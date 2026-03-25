from __future__ import annotations

import json
import re
import subprocess
import time
from pathlib import Path
from typing import Any

from docx import Document
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.shared import Pt
from pptx import Presentation


def markdown_to_docx(markdown_text: str, output_path: Path) -> None:
    document = Document()
    style = document.styles["Normal"]
    style.font.name = "PingFang SC"
    style.font.size = Pt(11)
    for raw_line in markdown_text.splitlines():
        line = raw_line.rstrip()
        if not line:
            document.add_paragraph("")
            continue
        if line.startswith("# "):
            paragraph = document.add_paragraph()
            paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            run = paragraph.add_run(line[2:].strip())
            run.bold = True
            run.font.size = Pt(16)
            continue
        if line.startswith("## "):
            document.add_heading(line[3:].strip(), level=1)
            continue
        if re.match(r"^\d+\.\s+", line):
            document.add_paragraph(re.sub(r"^\d+\.\s+", "", line), style="List Number")
            continue
        if line.startswith("- "):
            document.add_paragraph(line[2:].strip(), style="List Bullet")
            continue
        document.add_paragraph(line)
    document.save(str(output_path))


def extract_markdown_sections(markdown: str) -> dict[str, str]:
    current = "ROOT"
    sections: dict[str, list[str]] = {current: []}
    for line in markdown.splitlines():
        if line.startswith("#"):
            current = line.lstrip("#").strip()
            sections.setdefault(current, [])
            continue
        sections[current].append(line)
    return {key: "\n".join(lines).strip() for key, lines in sections.items()}


def compress_to_bullets(text: str, max_items: int = 4) -> list[str]:
    candidates = re.split(r"(?<=[。！？；])\s*", re.sub(r"\s+", " ", text))
    bullets = []
    for sentence in candidates:
        cleaned = sentence.strip(" -")
        if not cleaned:
            continue
        bullets.append(cleaned[:88] + ("..." if len(cleaned) > 88 else ""))
        if len(bullets) >= max_items:
            break
    return bullets or ["待根据正文补充。"]


def references_to_bullets(text: str, max_items: int = 4) -> list[str]:
    bullets = []
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        cleaned = re.sub(r"^\d+[.)、]\s*", "", line)
        bullets.append(cleaned[:88] + ("..." if len(cleaned) > 88 else ""))
        if len(bullets) >= max_items:
            break
    return bullets or compress_to_bullets(text, max_items=max_items)


def render_bundle(
    markdown: str,
    title: str,
    deck_outline: list[dict[str, Any]],
    output_dir: Path,
    basic_info: dict[str, str],
    footer_label: str,
) -> dict[str, str]:
    safe_title = re.sub(r"[/:]", "-", title).strip()
    report_md_path = output_dir / f"{safe_title}-开题报告.md"
    report_docx_path = output_dir / f"{safe_title}-开题报告.docx"
    deck_path = output_dir / f"{safe_title}-开题答辩PPT.pptx"
    notes_md_path = output_dir / f"{safe_title}-开题答辩讲稿.md"
    notes_docx_path = output_dir / f"{safe_title}-开题答辩讲稿.docx"
    report_md_path.write_text(markdown, encoding="utf-8")
    markdown_to_docx(markdown, report_docx_path)

    sections = extract_markdown_sections(markdown)
    payloads = _build_deck_payloads(
        title=title,
        deck_outline=deck_outline,
        basic_info=basic_info,
        sections=sections,
    )
    _render_presentation(
        payloads=payloads,
        title=title,
        basic_info=basic_info,
        footer_label=footer_label,
        output_path=deck_path,
    )
    notes_sections = _build_notes_sections(
        payloads=payloads,
        title=title,
        basic_info=basic_info,
    )
    _write_speaker_notes(deck_path, notes_sections)
    notes_markdown = _render_notes_markdown(notes_sections, title)
    notes_md_path.write_text(notes_markdown, encoding="utf-8")
    markdown_to_docx(notes_markdown, notes_docx_path)
    return {
        "report_markdown_path": str(report_md_path),
        "report_docx_path": str(report_docx_path),
        "deck_pptx_path": str(deck_path),
        "notes_md_path": str(notes_md_path),
        "notes_docx_path": str(notes_docx_path),
    }


def _build_deck_payloads(
    *,
    title: str,
    deck_outline: list[dict[str, Any]],
    basic_info: dict[str, str],
    sections: dict[str, str],
) -> list[dict[str, Any]]:
    section_map = {
        "background": compress_to_bullets(sections.get("选题的背景与问题的提出", "")),
        "problem": compress_to_bullets(sections.get("选题的背景与问题的提出", "")),
        "questions": compress_to_bullets(sections.get("研究设计", "")),
        "methods": compress_to_bullets(sections.get("研究设计", "")),
        "sources": references_to_bullets(sections.get("主要参考文献目录", "")),
        "outline": compress_to_bullets(sections.get("研究设计", "")),
        "timeline": compress_to_bullets(sections.get("研究方案及其进度安排", "")),
    }
    role_map = {
        "cover": "汇报开场",
        "background": "管理场景切入",
        "problem": "核心矛盾边界",
        "questions": "拟回答的问题",
        "methods": "方法与样本设计",
        "sources": "文献与资料条件",
        "outline": "章节与论证逻辑",
        "timeline": "里程碑与节奏",
    }
    summary_map = {
        "sources": "现有文献已覆盖需求管理、客户需求转译与跨部门协同等关键方向。",
        "timeline": "从开题、调研、分析到成稿，整体按阶段推进并预留修改缓冲。",
    }
    payloads = []
    for index, slide in enumerate(deck_outline, start=1):
        if slide["id"] == "cover":
            bullets = [
                f"题目：{title}",
                f"学生：{basic_info.get('student_name') or '待确认'}",
                f"导师：{basic_info.get('mentor_name') or '待确认'}",
                f"论文类型：{basic_info.get('thesis_type') or '专题研究类'}",
            ]
        else:
            bullets = section_map.get(slide["id"], ["待根据正文补充。"])
        payloads.append(
            {
                "page_no": index,
                "title": slide["title"],
                "id": slide["id"],
                "bullets": bullets[:4],
                "summary": summary_map.get(slide["id"]) or (bullets[0] if bullets else "待根据正文补充。"),
                "role": role_map.get(slide["id"], "研究说明"),
                "duration_seconds": slide["duration_seconds"],
            }
        )
    return payloads


def _render_presentation(
    *,
    payloads: list[dict[str, Any]],
    title: str,
    basic_info: dict[str, str],
    footer_label: str,
    output_path: Path,
) -> None:
    backend_root = Path(__file__).resolve().parents[2]
    renderer_path = backend_root / "app" / "ppt" / "render_deck.mjs"
    spec_path = output_path.with_suffix(".deck.json")
    previous_mtime = output_path.stat().st_mtime if output_path.exists() else None
    spec = {
        "title": title,
        "footerLabel": footer_label,
        "basicInfo": basic_info,
        "outputPath": str(output_path),
        "slides": payloads,
    }
    spec_path.write_text(
        json.dumps(spec, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    result = subprocess.run(
        ["node", str(renderer_path), str(spec_path)],
        cwd=str(backend_root),
        capture_output=True,
        text=True,
    )
    if result.returncode != 0:
        detail = (result.stderr or result.stdout or "").strip()
        raise RuntimeError(f"PPT 导出失败：{detail or 'render_deck.mjs exited with non-zero status'}")
    _wait_for_fresh_output(output_path, previous_mtime)


def _build_notes_sections(
    *,
    payloads: list[dict[str, Any]],
    title: str,
    basic_info: dict[str, str],
) -> list[dict[str, Any]]:
    return [
        {
            "page_no": payload["page_no"],
            "title": payload["title"],
            "duration_seconds": payload["duration_seconds"],
            "speech": _compose_speaker_script(payload, title, basic_info),
        }
        for payload in payloads
    ]


def _render_notes_markdown(notes_sections: list[dict[str, Any]], title: str) -> str:
    blocks = [f"# 《{title}》开题答辩讲稿", ""]
    for section in notes_sections:
        blocks.extend(
            [
                f"## 第{section['page_no']}页 {section['title']}",
                "",
                f"建议时长：{section['duration_seconds']}秒",
                "",
                "发言稿：",
                "",
                section["speech"],
                "",
            ]
        )
    return "\n".join(blocks).strip() + "\n"


def _compose_speaker_script(
    payload: dict[str, Any],
    report_title: str,
    basic_info: dict[str, str],
) -> str:
    if payload["id"] == "cover":
        student = basic_info.get("student_name") or "待确认"
        school = basic_info.get("school_name") or "待确认学校"
        mentor = basic_info.get("mentor_name") or "待确认导师"
        thesis_type = basic_info.get("thesis_type") or "专题研究类"
        return (
            f"各位老师好，我是{school}{thesis_type}学员{student}，"
            f"我的论文题目是《{report_title}》，导师是{mentor}。"
            "本次汇报会围绕选题背景、核心问题、研究设计、资料基础和进度安排展开，"
            "重点说明这个题目为什么值得研究、为什么现在能做，以及后续准备怎样把研究落到实处。"
        )
    if payload["id"] == "sources":
        references = [_reference_topic(item) for item in payload["bullets"] if str(item).strip()]
        examples = "、".join(references[:2]) if references else "现有核心文献"
        return (
            "这一页重点交代文献与资料基础，证明研究已经具备学术支撑和一手材料条件。"
            f"目前已经初步整理了{examples}等代表性研究，同时还会结合企业内部资料、访谈记录和项目复盘材料做交叉验证。"
            "这说明本研究后续写作具备比较扎实的文献与资料条件。"
        )

    opener_map = {
        "background": "这一页我先交代选题背景，说明这个研究不是凭空设想，而是从真实管理场景中提炼出来的。",
        "problem": "在背景明确以后，这一页进一步聚焦核心问题，回答到底要解决什么、问题边界在哪里。",
        "questions": "这一页主要说明研究目的和预期回答的关键问题，确保研究方向保持集中。",
        "methods": "接下来我说明研究设计，也就是准备怎样开展研究、采用什么方法、依靠什么资料完成论证。",
        "sources": "这一页重点交代文献与资料基础，证明研究已经具备学术支撑和一手材料条件。",
        "outline": "在研究方案层面，我会把论文结构和主要分析路径说明清楚，便于各位老师判断可执行性。",
        "timeline": "最后一页我汇报整体进度安排，说明开题之后各阶段准备如何推进。",
    }
    closer_map = {
        "background": "这些背景信息共同说明，研究对象是真实存在、而且已经具备研究紧迫性的管理问题。",
        "problem": "通过这样的界定，后续研究就能围绕清晰主线展开，避免把问题做散、做泛。",
        "questions": "这样处理的目的，是让论文始终围绕应用研究导向展开，而不是停留在概念讨论。",
        "methods": "总体上，这样的设计能够兼顾问题真实性、资料可得性和研究可操作性。",
        "sources": "这也意味着后续论文写作具备比较扎实的证据基础，而不是只依赖泛泛的二手材料。",
        "outline": "按这个框架推进，论文从问题提出到方案落地会形成比较完整的逻辑闭环。",
        "timeline": "我的汇报完毕，请各位老师批评指正。",
    }
    order_words = ("第一", "第二", "第三", "第四")
    speech_parts = [opener_map.get(payload["id"], f"这一页主要汇报{payload['title']}。")]
    normalized_bullets = [_normalize_speech_text(item) for item in payload["bullets"] if str(item).strip()]
    for index, bullet in enumerate(normalized_bullets):
        prefix = order_words[index] if index < len(order_words) else f"第{index + 1}"
        speech_parts.append(f"{prefix}，{bullet}")
    speech_parts.append(closer_map.get(payload["id"], "以上是本页的主要内容。"))
    return " ".join(part.strip() for part in speech_parts if part.strip())


def _normalize_speech_text(text: str) -> str:
    cleaned = re.sub(r"\s+", " ", str(text or "")).strip(" -")
    cleaned = cleaned.replace("...", "。")
    cleaned = cleaned.strip("。；，、 ")
    return cleaned + "。"


def _reference_topic(text: str) -> str:
    parts = [item.strip() for item in re.split(r"[.。]", str(text or "")) if item.strip()]
    if len(parts) >= 2:
        return parts[1]
    return parts[0] if parts else "相关研究"


def _write_speaker_notes(presentation_path: Path, notes_sections: list[dict[str, Any]]) -> None:
    presentation = Presentation(str(presentation_path))
    notes_by_page = {item["page_no"]: item for item in notes_sections}
    if len(presentation.slides) != len(notes_sections):
        raise ValueError("PPT 页数与讲稿页数不一致，无法写入备注")

    for index, slide in enumerate(presentation.slides, start=1):
        note = notes_by_page.get(index)
        if note is None:
            continue
        text_frame = slide.notes_slide.notes_text_frame
        if text_frame is None:
            raise ValueError(f"第{index}页缺少备注文本框，无法写入讲稿")
        text_frame.clear()
        text_frame.word_wrap = True
        parts = [
            f"建议时长：{note['duration_seconds']}秒",
            note["speech"],
        ]
        text_frame.paragraphs[0].text = parts[0]
        for part in parts[1:]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = part

    presentation.save(str(presentation_path))


def _wait_for_fresh_output(
    output_path: Path,
    previous_mtime: float | None,
    timeout_seconds: float = 5.0,
) -> None:
    deadline = time.time() + timeout_seconds
    while time.time() < deadline:
        if output_path.exists():
            stat = output_path.stat()
            if stat.st_size > 0 and (
                previous_mtime is None or stat.st_mtime > previous_mtime
            ):
                return
        time.sleep(0.1)
    raise RuntimeError(f"PPT 导出失败：未在 {timeout_seconds:.1f} 秒内确认文件落盘 {output_path}")

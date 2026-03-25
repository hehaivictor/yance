from __future__ import annotations

import json
import re
from collections import defaultdict
from datetime import date
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

import bibtexparser
from docx import Document
from pptx import Presentation

from .llm import LLMError, extract_image_text, is_enabled

try:
    from pypdf import PdfReader
except Exception:  # pragma: no cover
    PdfReader = None


SUPPORTED_TEXT_EXTENSIONS = {".md", ".txt", ".docx", ".pdf", ".pptx"}
SUPPORTED_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"}
SUPPORTED_PARSE_EXTENSIONS = SUPPORTED_TEXT_EXTENSIONS | SUPPORTED_IMAGE_EXTENSIONS

FIELD_HINT_PATTERNS = {
    "student_name": [r"学生[：:]\s*([^\s/，,。]{2,8})"],
    "student_id": [r"学号[：:]\s*([0-9]{8,20})"],
    "mentor_name": [
        r"指导教师[：:]\s*([^\s，,。]{2,12})",
        r"导师[：:]\s*([^\s，,。]{2,12})",
    ],
    "mentor_title": [r"(教授|副教授|研究员|副研究员|讲师)"],
    "research_direction": [r"研究方向[：:]\s*([^\n]+)"],
    "thesis_type": [r"论文类型[：:]\s*([^\n]+)"],
    "school_name": [r"(武汉大学)"],
    "program_name": [r"\b(EMBA|MBA|MPAcc)\b"],
}


def read_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".md", ".txt"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".docx":
        document = Document(str(path))
        return "\n".join(paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip())
    if suffix == ".pdf" and PdfReader:
        reader = PdfReader(str(path))
        chunks: list[str] = []
        for page in reader.pages[:10]:
            chunks.append(page.extract_text() or "")
        return "\n".join(chunks)
    if suffix == ".pptx":
        deck = Presentation(str(path))
        chunks: list[str] = []
        for slide in deck.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    chunks.append(text.strip())
        return "\n".join(item for item in chunks if item)
    if suffix in SUPPORTED_IMAGE_EXTENSIONS:
        if not is_enabled():
            return ""
        try:
            return _normalize_ocr_text(extract_image_text(path))
        except LLMError:
            return ""
    return ""


def extract_urls(text: str) -> list[str]:
    raw = re.findall(r"https?://[^\s)>\]]+", text)
    cleaned = [item.rstrip(".,;:)]}>。；，") for item in raw]
    return list(dict.fromkeys(cleaned))


def extract_metadata(text: str) -> dict[str, str]:
    values: dict[str, str] = {}
    for field, patterns in FIELD_HINT_PATTERNS.items():
        for pattern in patterns:
            match = re.search(pattern, text)
            if match:
                values[field] = match.group(1).strip("《》 ")
                break
    mentor = values.get("mentor_name")
    if mentor:
        values["mentor_name"] = re.sub(r"(教授|副教授|研究员|副研究员|讲师)$", "", mentor)
    return values


def classify_local_file(path: Path) -> str:
    name = path.name.lower()
    if any(key in name for key in ["写作指南", "通知", "登记表", "要求", "学位论文"]):
        return "school_requirement"
    if any(key in name for key in ["开题报告", "开题", "答辩", "讲稿", "初稿", "正式稿", "重写版"]):
        return "draft"
    if any(key in name for key in ["访谈", "方案", "需求", "项目", "复盘", "数据", "内部", "agent"]):
        return "internal_material"
    return "background_material"


def grade_for_local_category(category: str) -> str:
    if category in {"school_requirement", "internal_material"}:
        return "A"
    if category == "background_material":
        return "B"
    return "C"


def evidence_summary_from_text(text: str, limit: int = 240) -> str:
    compact = re.sub(r"\s+", " ", text).strip()
    return compact[:limit] + ("..." if len(compact) > limit else "")


def parse_ris(text: str) -> list[dict[str, Any]]:
    entries: list[dict[str, Any]] = []
    current: dict[str, list[str]] = defaultdict(list)
    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        if not line:
            continue
        if line.startswith("ER  -"):
            entries.append(_ris_entry_to_metadata(current))
            current = defaultdict(list)
            continue
        if "  -" not in line:
            continue
        tag, value = line.split("  -", 1)
        current[tag.strip()].append(value.strip())
    if current:
        entries.append(_ris_entry_to_metadata(current))
    return [entry for entry in entries if entry.get("title")]


def _ris_entry_to_metadata(entry: dict[str, list[str]]) -> dict[str, Any]:
    author_list = entry.get("AU") or entry.get("A1") or []
    title = (entry.get("TI") or entry.get("T1") or [""])[0]
    journal = (entry.get("JO") or entry.get("JF") or entry.get("T2") or [""])[0]
    year = ""
    if entry.get("PY"):
        year = entry["PY"][0][:4]
    elif entry.get("Y1"):
        year = entry["Y1"][0][:4]
    doi = (entry.get("DO") or [""])[0]
    url = (entry.get("UR") or [""])[0]
    return {
        "author": "; ".join(author_list),
        "title": title,
        "year": year,
        "source": journal,
        "doi": doi,
        "url": url,
    }


def parse_bibtex(text: str) -> list[dict[str, Any]]:
    database = bibtexparser.loads(text)
    entries: list[dict[str, Any]] = []
    for entry in database.entries:
        entries.append(
            {
                "author": entry.get("author", ""),
                "title": entry.get("title", "").strip("{}"),
                "year": entry.get("year", ""),
                "source": entry.get("journal") or entry.get("booktitle") or entry.get("publisher", ""),
                "doi": entry.get("doi", ""),
                "url": entry.get("url", ""),
            }
        )
    return [entry for entry in entries if entry.get("title")]


def citation_completeness_score(metadata: dict[str, Any]) -> tuple[bool, list[str]]:
    missing = []
    for key in ["author", "title", "year", "source"]:
        if not str(metadata.get(key, "")).strip():
            missing.append(key)
    has_locator = bool(str(metadata.get("doi", "")).strip() or str(metadata.get("url", "")).strip())
    if not has_locator:
        missing.append("doi/url")
    return not missing, missing


def build_citation_reference(metadata: dict[str, Any], index: int) -> str:
    author = metadata.get("author") or "作者待补"
    year = metadata.get("year") or "年份待补"
    title = metadata.get("title") or "题名待补"
    source = metadata.get("source") or "来源待补"
    doi = metadata.get("doi") or metadata.get("url") or ""
    tail = f". {doi}" if doi else ""
    return f"[{index}] {author}. {title}[J/OL]. {source}, {year}{tail}"


def parse_citation_file(path: Path) -> list[dict[str, Any]]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    suffix = path.suffix.lower()
    if suffix == ".ris":
        return parse_ris(text)
    if suffix == ".bib":
        return parse_bibtex(text)
    if suffix == ".json":
        payload = json.loads(text)
        if isinstance(payload, list):
            return [item for item in payload if isinstance(item, dict)]
    return []


def extract_page_date(text: str) -> str:
    patterns = [
        r"(\d{4}-\d{2}-\d{2})",
        r"(\d{4}/\d{2}/\d{2})",
        r"(\d{4}年\d{1,2}月\d{1,2}日)",
    ]
    for pattern in patterns:
        match = re.search(pattern, text)
        if match:
            return match.group(1)
    return ""


def domain_from_url(url: str) -> str:
    return urlparse(url).netloc.lower()


def slug_from_title(title: str) -> str:
    cleaned = re.sub(r"[^\w\u4e00-\u9fff]+", "-", title).strip("-")
    return cleaned[:60] or "workspace"


def today_string() -> str:
    return date.today().isoformat()


def _normalize_ocr_text(text: str) -> str:
    cleaned = str(text or "").strip()
    if cleaned.startswith("```"):
        cleaned = cleaned.strip("`").strip()
        if cleaned.lower().startswith("text"):
            cleaned = cleaned[4:].strip()
    if cleaned in {"空字符串", "未识别到文字", "未检测到文字", "无可辨识文字"}:
        return ""
    return cleaned
